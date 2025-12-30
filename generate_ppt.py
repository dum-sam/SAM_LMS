from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()

    # Define Slide Layouts
    TITLE_SLIDE_LAYOUT = prs.slide_layouts[0]
    BULLET_SLIDE_LAYOUT = prs.slide_layouts[1]

    # --- Slide 1: Title ---
    slide = prs.slides.add_slide(TITLE_SLIDE_LAYOUT)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "SAM-LMS: Master the Future"
    subtitle.text = "Strategic Advanced Method E-Learning\nFinal Year Project Presentation"

    # Helper function to add slides
    def add_slide(title_text, content_points):
        slide = prs.slides.add_slide(BULLET_SLIDE_LAYOUT)
        title = slide.shapes.title
        title.text = title_text
        
        tf = slide.placeholders[1].text_frame
        tf.clear()  # Clear default paragraph
        
        for point in content_points:
            p = tf.add_paragraph()
            p.text = point
            p.level = 0
            if isinstance(point, tuple): # Handle sub-points if needed in future
                 pass 

    # --- Slide 2: Introduction ---
    add_slide("1. Introduction", [
        "Project Name: SAM-LMS (Strategic Advanced Method E-Learning)",
        "Tagline: More than just videos. A structured career roadmap.",
        "Problem: Traditional LMS platforms are static, boring, and lack motivation.",
        "Solution: A Modern, Gamified, and Path-Driven learning platform.",
        "Core Value: Groups courses into Career Paths and tracks real progression."
    ])

    # --- Slide 3: Architecture ---
    add_slide("2. Unseen Architecture (How it Works)", [
        "Fully distributed cloud architecture (Not just Localhost).",
        "Frontend: HTML5, Bootstrap 5, Vanilla JS (Fast & Responsive).",
        "Backend: Python 3.12 + Django 5.0 (Secure & Scalable).",
        "Database: PostgreSQL (Hosted on Neon/Vercel).",
        "Media: Cloudinary CDN for fast video streaming.",
        "Real-time: AJAX Polling for the Community Chat modules."
    ])

    # --- Slide 4: Key Feature - Learning Paths ---
    add_slide("3. Learning Paths (Career Tracks)", [
        "We don't just dump courses on the user.",
        "Curated Paths: Beginner -> Intermediate -> Expert.",
        "Example: 'Full Stack Architect' or 'Data Scientist'.",
        "Real Tracking: Persistent progress bars per path.",
        "Goal: Guides students accurately through their career journey."
    ])

    # --- Slide 5: Key Feature - Gamification ---
    add_slide("4. Gamification & Motivation", [
        "Problem: High student dropout rates.",
        "Solution: 'Video Game' style progression system.",
        "XP System: Daily Check-ins award +10 XP.",
        "Leveling: Level = (XP / 100) + 1.",
        "Badges: Visual rewards (e.g., 'Level 5') displayed on dashboards."
    ])

    # --- Slide 6: Key Feature - Community ---
    add_slide("5. The Community", [
        "Learning is social.",
        "Integrated Discord-like chat system inside the dashboard.",
        "Real-time updates using efficient polling.",
        "Role-Based: Instructors have distinct visual highlights.",
        "Secure: Only enrolled students can access discussions."
    ])

    # --- Slide 7: Security & Roles ---
    add_slide("6. Security Measures", [
        "CSRF Protection: Prevents cross-site attacks.",
        "RBAC (Role-Based Access Control):",
        " - Student: Read content, Write chat/notes.",
        " - Instructor: Create/Manage courses.",
        " - Admin: Full system control.",
        "Encrypted Sessions and Password Hashing (PBKDF2)."
    ])

    # --- Slide 8: Deployment Strategy ---
    add_slide("7. Deployment (Production Ready)", [
        "Hosted on Vercel (Serverless Architecture).",
        "Database on Neon (Serverless Postgres).",
        "Media on Cloudinary.",
        "CI/CD: GitHub integration.",
        " - Every push to 'master' deploys live in ~30 seconds.",
        " - No manual server restarts required."
    ])

    # --- Slide 9: Dashboard Walkthrough ---
    add_slide("8. The Dashboard", [
        "Active Paths: Shows major career goals at a glance.",
        "Smart Recommendations: Suggests courses based on history.",
        "Stats Card: Visual feedback (Active, Completed, Certs).",
        "Profile: Displays User Level, XP, and Roles."
    ])

    # --- Slide 10: Future & Conclusion ---
    add_slide("9. Future & Conclusion", [
        "Future Enhancements:",
        " - AI Tutor (OpenAI API integration).",
        " - Mobile App (React Native via Django REST API).",
        " - Payments (Stripe/Razorpay).",
        "Conclusion: SAM-LMS is a scalable, educational ecosystem.",
        "Thank You!"
    ])

    output_path = "SAMLMS_Presentation.pptx"
    prs.save(output_path)
    print(f"Successfully saved {output_path}")

if __name__ == "__main__":
    create_presentation()
